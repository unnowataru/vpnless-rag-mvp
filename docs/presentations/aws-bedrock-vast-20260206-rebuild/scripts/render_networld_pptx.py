#!/usr/bin/env python3
import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Pt

try:
    from PIL import Image
except Exception:
    Image = None


LAYOUT_RULES = {
    "cover": {
        "names": ["表紙-2"],
        "fallback_index": 0,
        "title_idx": 0,
    },
    "section": {
        "names": ["表紙-1"],
        "fallback_index": 1,
        "title_idx": 0,
    },
    "content": {
        "names": ["Page-1"],
        "fallback_index": 2,
        "title_idx": 0,
        "body_idx": 10,
    },
    "content_alt": {
        "names": ["タイトルとコンテンツ"],
        "fallback_index": 4,
        "title_idx": 0,
        "body_idx": 1,
    },
}


def _safe_text(value, fallback=""):
    if value is None:
        return fallback
    text = str(value).strip()
    return text if text else fallback


def _to_string_list(value):
    if not isinstance(value, list):
        return []
    out = []
    for item in value:
        text = str(item).strip()
        if text:
            out.append(text)
    return out


def _clear_existing_slides(prs):
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide_id = prs.slides._sldIdLst[idx]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[idx]


def _find_layout(prs, layout_key):
    rule = LAYOUT_RULES[layout_key]
    for name in rule["names"]:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
    fallback = rule["fallback_index"]
    if fallback < len(prs.slide_layouts):
        return prs.slide_layouts[fallback]
    raise ValueError(f"Layout not found for key: {layout_key}")


def _find_placeholder(slide, idx):
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == idx:
            return shp
    return None


def _set_title(slide, idx, text):
    placeholder = _find_placeholder(slide, idx)
    if placeholder is not None and placeholder.has_text_frame:
        placeholder.text = text


def _set_body(slide, idx, body, bullets):
    placeholder = _find_placeholder(slide, idx)
    if placeholder is None or not placeholder.has_text_frame:
        return
    frame = placeholder.text_frame
    frame.clear()
    if bullets:
        frame.text = bullets[0]
        for bullet in bullets[1:]:
            para = frame.add_paragraph()
            para.text = bullet
            para.level = 0
    else:
        frame.text = body


def _clear_body_placeholder_text(slide, idx):
    placeholder = _find_placeholder(slide, idx)
    if placeholder is not None and placeholder.has_text_frame:
        placeholder.text_frame.clear()


def _add_picture_fit(slide, image_path, left, top, width, height):
    if Image is None:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    if img_w <= 0 or img_h <= 0:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return

    box_w = int(width)
    box_h = int(height)
    scale = min(box_w / img_w, box_h / img_h)
    draw_w = int(img_w * scale)
    draw_h = int(img_h * scale)

    draw_left = Emu(int(left) + (box_w - draw_w) // 2)
    draw_top = Emu(int(top) + (box_h - draw_h) // 2)
    slide.shapes.add_picture(str(image_path), draw_left, draw_top, width=Emu(draw_w), height=Emu(draw_h))


def _add_caption(slide, caption, left, top, width, height):
    if not caption:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.text = caption
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(80, 80, 80)


def _render_visual(slide, visual, slide_width, slide_height):
    if not isinstance(visual, dict):
        return
    visual_type = _safe_text(visual.get("type"), "none")
    if visual_type == "none":
        return

    image_path = Path(_safe_text(visual.get("image_path"), ""))
    if not image_path.exists():
        return

    margin_x = Emu(420000)
    top = Emu(1300000)
    usable_h = Emu(int(slide_height) - 1900000)
    usable_w = Emu(int(slide_width) - 2 * int(margin_x))

    if visual_type == "image_with_caption":
        caption_h = Emu(280000)
        image_h = Emu(int(usable_h) - int(caption_h) - 50000)
        _add_picture_fit(slide, image_path, margin_x, top, usable_w, image_h)
        _add_caption(slide, _safe_text(visual.get("caption"), ""), margin_x, Emu(int(top) + int(image_h) + 30000), usable_w, caption_h)
        return

    if visual_type == "image_with_callouts":
        left_w = Emu(int(usable_w) * 66 // 100)
        right_gap = Emu(120000)
        right_left = Emu(int(margin_x) + int(left_w) + int(right_gap))
        right_w = Emu(int(usable_w) - int(left_w) - int(right_gap))

        caption_h = Emu(260000)
        image_h = Emu(int(usable_h) - int(caption_h) - 50000)
        _add_picture_fit(slide, image_path, margin_x, top, left_w, image_h)
        _add_caption(slide, _safe_text(visual.get("caption"), ""), margin_x, Emu(int(top) + int(image_h) + 30000), left_w, caption_h)

        callouts = _to_string_list(visual.get("callouts"))
        if not callouts:
            return

        box_gap = Emu(110000)
        total_gap = int(box_gap) * (len(callouts) - 1)
        box_h = max(Emu(300000), Emu((int(image_h) - total_gap) // len(callouts)))

        cursor_top = top
        for text in callouts:
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                right_left,
                cursor_top,
                right_w,
                box_h,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
            shape.line.color.rgb = RGBColor(200, 205, 215)
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.text = text
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(40, 45, 55)
            cursor_top = Emu(int(cursor_top) + int(box_h) + int(box_gap))


def render(template_path, spec, output_path):
    prs = Presentation(str(template_path))
    _clear_existing_slides(prs)

    title = _safe_text(spec.get("deck_title"), "Untitled Deck")
    subtitle = _safe_text(spec.get("deck_subtitle"), "")

    cover_layout = _find_layout(prs, "cover")
    cover = prs.slides.add_slide(cover_layout)
    cover_text = title if not subtitle else f"{title}\n{subtitle}"
    _set_title(cover, LAYOUT_RULES["cover"]["title_idx"], cover_text)

    slides = spec.get("slides") if isinstance(spec.get("slides"), list) else []

    for idx, slide_spec in enumerate(slides, start=1):
        if not isinstance(slide_spec, dict):
            continue

        layout_key = _safe_text(slide_spec.get("layout"), "content")
        if layout_key not in ("section", "content", "content_alt"):
            layout_key = "content"

        slide = prs.slides.add_slide(_find_layout(prs, layout_key))
        slide_title = _safe_text(slide_spec.get("title"), f"Slide {idx}")
        _set_title(slide, LAYOUT_RULES[layout_key]["title_idx"], slide_title)

        body = _safe_text(slide_spec.get("body"), "")
        bullets = _to_string_list(slide_spec.get("bullets"))
        visual = slide_spec.get("visual")

        if layout_key in ("content", "content_alt"):
            body_idx = LAYOUT_RULES[layout_key]["body_idx"]
            if isinstance(visual, dict) and _safe_text(visual.get("type"), "none") != "none":
                _clear_body_placeholder_text(slide, body_idx)
                _render_visual(slide, visual, prs.slide_width, prs.slide_height)
            else:
                _set_body(slide, body_idx, body, bullets)

        notes = _safe_text(slide_spec.get("notes"), "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Render Networld template-based PPTX from slide spec.")
    parser.add_argument("--template", required=True, help="Path to Networld-Basic.pptx")
    parser.add_argument("--input", required=True, help="Path to slide-spec JSON")
    parser.add_argument("--output", required=True, help="Path to output .pptx")
    args = parser.parse_args()

    template_path = Path(args.template)
    input_path = Path(args.input)
    output_path = Path(args.output)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    if not input_path.exists():
        raise FileNotFoundError(f"Input JSON not found: {input_path}")

    spec = json.loads(input_path.read_text(encoding="utf-8"))
    render(template_path, spec, output_path)
    print(str(output_path))


if __name__ == "__main__":
    main()
